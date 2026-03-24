from pathlib import Path
import pandas as pd
from pypdf import PdfReader
from docx import Document
from pptx import Presentation

def load_pdf(path: Path) -> str:
    reader = PdfReader(str(path))
    texts = []
    for page_num, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if text.strip():
            texts.append(f"[Page {page_num}]\n{text}")
    return "\n".join(texts)

def load_docx(path: Path) -> str:
    doc = Document(str(path))
    return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])

def load_pptx(path: Path) -> str:
    prs = Presentation(str(path))
    slides_text = []
    for idx, slide in enumerate(prs.slides, start=1):
        text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if shape.text.strip():
                    text_parts.append(shape.text)
        if text_parts:
            slides_text.append(f"[Slide {idx}]\n" + "\n".join(text_parts))
    return "\n".join(slides_text)

def load_xlsx(path: Path) -> str:
    xls = pd.ExcelFile(path)
    sheet_texts = []
    for sheet_name in xls.sheet_names:
        df = pd.read_excel(path, sheet_name=sheet_name)
        content = df.fillna("").astype(str).to_string(index=False)
        sheet_texts.append(f"[Sheet: {sheet_name}]\n{content}")
    return "\n".join(sheet_texts)

def load_file(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return load_pdf(path)
    elif suffix == ".docx":
        return load_docx(path)
    elif suffix == ".pptx":
        return load_pptx(path)
    elif suffix in [".xlsx", ".xlsm"]:
        return load_xlsx(path)
    else:
        return ""